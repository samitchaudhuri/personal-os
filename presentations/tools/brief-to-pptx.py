#!/usr/bin/env python3
"""Build a .pptx from a generated PPT brief.

Deterministic local alternative to the Claude PowerPoint connector. Reads the
brief produced by marp-to-ppt-brief.js and renders one slide per "## Slide N"
section, placing each diagram full-width below the body text (the layout the
connector got wrong). Diagrams are embedded as PNGs (python-pptx cannot embed
SVG), so render PNGs first and pass their directory.

Usage:
  brief-to-pptx.py "<brief.md>" "<png-dir>" "<out.pptx>" ["<palette.json>"]

Layout (margins, gaps, diagram cap, PPTX type spacing) comes from
``themes/oram-common.json`` via ``load_layout()``. The optional palette JSON
themes slide chrome — colors, font, chrome toggles, and point sizes — parallel
to the Mermaid theme that themes the diagrams. Recognized palette keys: font,
background, ink, accent, takeawayFill, takeawayText, showKicker, showSlideNumber,
showTakeawayBand, and sizes (title/kicker/body/heading/takeaway/number). Missing
keys fall back to the built-in defaults below, so older 4-color palettes still work.
"""

import json
import re
import struct
import sys
from pathlib import Path

from reportlab.pdfbase.pdfmetrics import stringWidth
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

_DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DRAWING_TAG = f"{{{_DRAWING_NS}}}"  # Clark notation for lxml element find()
_TABLE_BORDER_TAG = {"left": "lnL", "right": "lnR", "top": "lnT", "bottom": "lnB"}
_TABLE_TBL_BORDER_SIDES = ("left", "right", "top", "bottom", "insideH", "insideV")
TABLE_BORDER_WIDTH_PT = 1.0

THEMES_DIR = Path(__file__).resolve().parent.parent / "themes"
DEFAULT_COMMON_PATH = THEMES_DIR / "oram-common.json"

# Palette defaults (light ORAM). Colors/font/toggles/sizes overridden by palette JSON.
ACCENT = RGBColor(0x4A, 0x6F, 0xA5)
INK = RGBColor(0x1A, 0x1A, 0x1A)
TAKEAWAY_FILL = RGBColor(0xE8, 0xEE, 0xF6)
TAKEAWAY_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF)
FONT = None  # None → PowerPoint default; a theme sets e.g. "Helvetica Neue".
SHOW_KICKER = True
SHOW_SLIDE_NUMBER = True
SHOW_TAKEAWAY_BAND = True
SIZES = {"title": 28, "kicker": 13, "body": 16, "heading": 17, "takeaway": 15, "number": 13}
TABLE_STRIPE = "row"
TABLE_HEADER_FILL = RGBColor(0xC9, 0xA8, 0x8E)
TABLE_HEADER_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TABLE_BAND_FILL = RGBColor(0xF5, 0xEB, 0xE4)
TABLE_CELL_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
VALID_TABLE_STRIPES = frozenset({"row", "column", "none"})

# Layout defaults mirror themes/oram-common.json until load_layout() runs.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN
CHROME_TOP = Inches(0.4)
KICKER_H = Inches(0.2)
GAP_AFTER_KICKER = Inches(0.06)
GAP_AFTER_TITLE = Inches(0.05)
GAP_AFTER_DIAGRAM = Inches(0.08)
TAKEAWAY_BOTTOM = Inches(0.35)
TAKEAWAY_BAND_H = Inches(0.85)
TAKEAWAY_TEXT_INSET = Inches(0.2)
MAX_DIAGRAM_H = Inches(260 / 96)
DIAGRAM_MAX_ZONE_FRACTION = 0.42
BODY_HEIGHT_PAD = Inches(0.2)
BODY_LINE_SPACING = 1.15
TITLE_LINE_HEIGHT = 1.3
TITLE_BOX_PAD = Inches(0.04)


def load_layout(common_path=None):
    """Load variant-agnostic layout from themes/oram-common.json."""
    global SLIDE_W, SLIDE_H, MARGIN, CONTENT_W, CHROME_TOP, KICKER_H
    global GAP_AFTER_KICKER, GAP_AFTER_TITLE, GAP_AFTER_DIAGRAM, TAKEAWAY_BOTTOM
    global TAKEAWAY_BAND_H, TAKEAWAY_TEXT_INSET, MAX_DIAGRAM_H, BODY_LINE_SPACING
    global DIAGRAM_MAX_ZONE_FRACTION, BODY_HEIGHT_PAD, TITLE_LINE_HEIGHT, TITLE_BOX_PAD

    path = Path(common_path) if common_path else DEFAULT_COMMON_PATH
    if not path.is_file():
        fail(f"ORAM common theme not found: {path}")

    data = json.loads(path.read_text(encoding="utf-8"))
    slide = data["slide"]
    layout = data["layout"]
    pptx = data["typography"]["pptx"]

    SLIDE_W = Inches(slide["widthIn"])
    SLIDE_H = Inches(slide["heightIn"])
    MARGIN = Inches(layout["marginIn"])
    CONTENT_W = SLIDE_W - 2 * MARGIN
    CHROME_TOP = Inches(layout["chromeTopIn"])
    KICKER_H = Inches(layout["kickerHeightIn"])
    GAP_AFTER_KICKER = Inches(layout["gapAfterKickerIn"])
    GAP_AFTER_TITLE = Inches(layout["gapAfterTitleIn"])
    GAP_AFTER_DIAGRAM = Inches(layout["gapAfterDiagramIn"])
    TAKEAWAY_BOTTOM = Inches(layout["takeawayBottomIn"])
    TAKEAWAY_BAND_H = Inches(layout["takeawayBandHeightIn"])
    MAX_DIAGRAM_H = Inches(layout["diagramMaxHeightPx"] / 96)
    BODY_LINE_SPACING = float(pptx["bodyLineSpacing"])
    TAKEAWAY_TEXT_INSET = Inches(pptx["takeawayTextInsetIn"])
    DIAGRAM_MAX_ZONE_FRACTION = float(layout.get("diagramMaxZoneFraction", 0.42))
    BODY_HEIGHT_PAD = Inches(layout.get("bodyHeightPadIn", 0.2))
    TITLE_LINE_HEIGHT = float(layout.get("titleLineHeight", 1.3))
    TITLE_BOX_PAD = Inches(layout.get("titleBoxPadIn", 0.04))


def load_theme(path):
    """Override the module-level chrome globals from a theme JSON file."""
    global ACCENT, INK, TAKEAWAY_FILL, TAKEAWAY_TEXT, BACKGROUND, FONT
    global SHOW_KICKER, SHOW_SLIDE_NUMBER, SHOW_TAKEAWAY_BAND, SIZES
    global TABLE_STRIPE, TABLE_HEADER_FILL, TABLE_HEADER_TEXT, TABLE_BAND_FILL
    global TABLE_CELL_TEXT
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    colors = {
        "ink": "INK",
        "accent": "ACCENT",
        "takeawayFill": "TAKEAWAY_FILL",
        "takeawayText": "TAKEAWAY_TEXT",
        "background": "BACKGROUND",
    }
    for key, name in colors.items():
        if data.get(key):
            globals()[name] = RGBColor.from_string(str(data[key]).lstrip("#"))
    table_colors = {
        "tableHeaderFill": "TABLE_HEADER_FILL",
        "tableHeaderText": "TABLE_HEADER_TEXT",
        "tableBandFill": "TABLE_BAND_FILL",
        "tableCellText": "TABLE_CELL_TEXT",
    }
    for key, name in table_colors.items():
        if data.get(key):
            globals()[name] = RGBColor.from_string(str(data[key]).lstrip("#"))
    # Legacy palette keys (bandFillB / bandFillA from older builds).
    if not data.get("tableBandFill") and data.get("tableBandFillB"):
        globals()["TABLE_BAND_FILL"] = RGBColor.from_string(
            str(data["tableBandFillB"]).lstrip("#")
        )
    if data.get("tableStripe"):
        stripe = str(data["tableStripe"]).strip().lower()
        if stripe not in VALID_TABLE_STRIPES:
            fail(
                f'Invalid tableStripe in {path}: "{data["tableStripe"]}" '
                f"(expected row, column, or none)."
            )
        globals()["TABLE_STRIPE"] = stripe
    if data.get("font"):
        FONT = str(data["font"])
    toggles = {
        "showKicker": "SHOW_KICKER",
        "showSlideNumber": "SHOW_SLIDE_NUMBER",
        "showTakeawayBand": "SHOW_TAKEAWAY_BAND",
    }
    for key, name in toggles.items():
        if key in data:
            globals()[name] = bool(data[key])
    sizes = data.get("sizes")
    if isinstance(sizes, dict):
        SIZES = {**SIZES, **{k: int(v) for k, v in sizes.items() if v}}


def _emu_inches(emu):
    return emu / 914400


def _content_width_pt():
    return _emu_inches(CONTENT_W) * 72


def _title_metrics_font():
    """ReportLab AFM face for bold title measurement (Helvetica Neue proxy)."""
    return "Helvetica-Bold"


def _title_line_count(title, font_pt, max_width_pt, metrics_font):
    """Count wrapped lines using font metrics (not character heuristics)."""
    text = title.strip()
    if not text:
        return 1

    lines = 1
    current = ""
    for word in text.split():
        candidate = f"{current} {word}".strip() if current else word
        if stringWidth(candidate, metrics_font, font_pt) <= max_width_pt:
            current = candidate
            continue
        if current:
            lines += 1
            current = word
        else:
            lines += 1
            current = ""
    return lines


def _title_box_height(title):
    """Estimate title textbox height from wrapped line count at the theme title size."""
    font_pt = SIZES["title"]
    line_h_in = font_pt / 72 * TITLE_LINE_HEIGHT
    lines = _title_line_count(
        title,
        font_pt,
        _content_width_pt(),
        _title_metrics_font(),
    )
    return Inches(line_h_in * lines + _emu_inches(TITLE_BOX_PAD))


def _body_metrics_font():
    return "Helvetica"


def _plain_text_for_measure(text):
    """Strip inline markdown so wrap measurement matches rendered text."""
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def _body_box_height(body):
    """Estimate body textbox height so diagram sizing can use leftover zone space."""
    if not body or not body.strip():
        return 0

    body_pt = SIZES["body"]
    heading_pt = SIZES["heading"]
    max_w = _content_width_pt()
    bullet_w = max_w - 14  # "•  " prefix in set_body

    height_pt = 0.0
    lines = [raw.rstrip() for raw in body.split("\n") if raw.strip()]
    for i, line in enumerate(lines):
        is_last = i == len(lines) - 1
        heading = re.match(r"^#{1,6}\s+(.*)$", line)
        if heading:
            wrapped = _title_line_count(
                _plain_text_for_measure(heading.group(1)),
                heading_pt,
                max_w,
                _title_metrics_font(),
            )
            height_pt += 4 + wrapped * heading_pt * BODY_LINE_SPACING + 3
        elif line.startswith("- "):
            wrapped = _title_line_count(
                _plain_text_for_measure(line[2:]),
                body_pt,
                bullet_w,
                _body_metrics_font(),
            )
            height_pt += wrapped * body_pt * BODY_LINE_SPACING + (2 if not is_last else 3)
        else:
            wrapped = _title_line_count(
                _plain_text_for_measure(line),
                body_pt,
                max_w,
                _body_metrics_font(),
            )
            height_pt += wrapped * body_pt * BODY_LINE_SPACING + (3 if is_last else 2)

    return Inches(height_pt / 72 + 0.04)


def _body_reserve(body):
    """Body zone height including safety pad so text does not spill into takeaway."""
    if not body or not body.strip():
        return 0
    return int(_body_box_height(body)) + int(BODY_HEIGHT_PAD)


def _diagram_max_height(zone_h, body=None):
    """Cap diagram height from zone budget; grow into space left by short body text."""
    zone = int(zone_h)
    abs_max = int(MAX_DIAGRAM_H)
    if body and body.strip():
        reserved = _body_reserve(body) + int(GAP_AFTER_DIAGRAM)
        available = max(0, zone - reserved)
        return min(available, abs_max)
    frac_cap = int(zone * DIAGRAM_MAX_ZONE_FRACTION)
    return min(frac_cap, abs_max)


def fail(msg):
    sys.stderr.write(f"Error: {msg}\n")
    sys.exit(1)


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        raise ValueError(f"Not a PNG: {path}")
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def parse_brief(text):
    # Drop everything before the first slide section.
    idx = text.find("\n## Slide ")
    if idx == -1:
        fail("No '## Slide N' sections found in brief.")
    body = text[idx:]
    chunks = re.split(r"\n## Slide \d+\s*\n", body)
    chunks = [c for c in chunks if c.strip()]
    return [parse_slide(c) for c in chunks]


def _field(chunk, label):
    m = re.search(rf"^\*\*{label}:\*\*\s*(.+)$", chunk, re.MULTILINE)
    return m.group(1).strip() if m else None


def _block(chunk, label):
    # Capture a multi-line block starting at **Label:** until the next **Field:** or end.
    m = re.search(
        rf"^\*\*{label}:\*\*\s*\n(.*?)(?=\n\*\*[A-Z][a-z]+:\*\*|\Z)",
        chunk,
        re.MULTILINE | re.DOTALL,
    )
    return m.group(1).strip() if m else None


def parse_slide(chunk):
    slide = {
        "kicker": _field(chunk, "Kicker"),
        "title": _field(chunk, "Title"),
        "takeaway": _field(chunk, "Takeaway"),
        "body": _block(chunk, "Body"),
        "notes": _block(chunk, "Notes"),
        "diagram": None,
    }
    dm = re.search(r"^\*\*Diagram:\*\*\s*`([^`]+)`", chunk, re.MULTILINE)
    if dm:
        slide["diagram"] = Path(dm.group(1)).stem  # strip .svg
    return slide


def add_rich_text(paragraph, text, size=None, color=None):
    """Render inline **bold** markup into runs."""
    for i, part in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.bold = i % 2 == 1
        if FONT:
            run.font.name = FONT
        if size is not None:
            run.font.size = size
        if color is not None:
            run.font.color.rgb = color


def set_body(tf, body):
    """Render brief body markdown into a fixed-height text frame.

    Marp PDF uses tight list spacing (``li { margin: 0.025em }``). PowerPoint's
    default paragraph spacing is much looser; without explicit line spacing and
    minimal ``space_after`` on bullets, body text overflows the vertical budget.
    """
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)

    lines = [raw.rstrip() for raw in body.split("\n") if raw.strip()]
    first = True
    for i, line in enumerate(lines):
        heading = re.match(r"^#{1,6}\s+(.*)$", line)
        bullet = line.startswith("- ")
        is_last = i == len(lines) - 1
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.level = 0
        para.line_spacing = BODY_LINE_SPACING
        if heading:
            para.space_before = Pt(4)
            para.space_after = Pt(3)
            add_rich_text(para, heading.group(1), size=Pt(SIZES["heading"]), color=INK)
            for run in para.runs:
                run.font.bold = True
        elif bullet:
            para.space_before = Pt(0)
            para.space_after = Pt(2 if not is_last else 3)
            add_rich_text(para, line[2:], size=Pt(SIZES["body"]), color=INK)
            if para.runs:
                # Visual bullet via prefix (keeps layout deterministic across renderers).
                para.runs[0].text = "•  " + para.runs[0].text
        else:
            para.space_before = Pt(0)
            para.space_after = Pt(3 if is_last else 2)
            add_rich_text(para, line, size=Pt(SIZES["body"]), color=INK)


def pick_blank_layout(prs):
    """Choose the best blank-like layout for generated slides.

    A fresh ``Presentation()`` has a blank at index 6 with only footer/date
    placeholders. Externally-authored bases name a layout ``Blank`` but still attach
    a slide-number placeholder — and index 6 is *not* blank (often a Title
    layout). Prefer an explicitly named Blank layout, then any layout without
    title placeholders, before falling back to index 6.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER

    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

    for layout in prs.slide_layouts:
        if layout.name.strip().lower() == "blank":
            return layout

    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout

    without_title = [
        layout
        for layout in prs.slide_layouts
        if not any(ph.placeholder_format.type in title_types for ph in layout.placeholders)
    ]
    if without_title:
        return min(without_title, key=lambda layout: len(layout.placeholders))

    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def clear_slide_placeholders(slide):
    """Remove inherited layout placeholders so generated chrome is not obscured."""
    for ph in list(slide.placeholders):
        el = ph.element
        el.getparent().remove(el)


def find_slide_number_shape(slide):
    """Return the top-right chrome textbox that holds the slide number.

    Assumes PowerPoint source slides were hand-fitted with the same number placement as
    ``render_slide`` (numeric text, upper-right band). Picks the rightmost match.
    """
    top_max = int(SLIDE_H * 0.18)
    right_min = int(SLIDE_W * 0.75)
    best = None
    best_left = -1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not re.fullmatch(r"\d{1,3}", text):
            continue
        if shape.top > top_max:
            continue
        if shape.left + shape.width < right_min:
            continue
        if shape.left > best_left:
            best_left = shape.left
            best = shape
    return best


def _clear_paragraph_runs(para):
    for run in list(para.runs):
        para._p.remove(run._r)


def _fill_slide_number_text_frame(tf, number):
    """Write slide number with theme accent (shared by render and stamp)."""
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    _clear_paragraph_runs(para)
    run = para.add_run()
    run.text = str(number)
    run.font.size = Pt(SIZES["number"])
    run.font.bold = True
    run.font.color.rgb = ACCENT
    if FONT:
        run.font.name = FONT
    for extra in tf.paragraphs[1:]:
        for run in extra.runs:
            run.text = ""


def _fill_kicker_text_frame(tf, role):
    """Write narrative role / kicker with theme accent (shared by render and stamp)."""
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    para = tf.paragraphs[0]
    _clear_paragraph_runs(para)
    run = para.add_run()
    run.text = role.upper()
    run.font.size = Pt(SIZES["kicker"])
    run.font.bold = True
    run.font.color.rgb = ACCENT
    if FONT:
        run.font.name = FONT


def _fill_title_text_frame(tf, title):
    """Write slide title with theme ink (shared by render and stamp)."""
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    para = tf.paragraphs[0]
    _clear_paragraph_runs(para)
    run = para.add_run()
    run.text = title
    run.font.size = Pt(SIZES["title"])
    run.font.bold = True
    run.font.color.rgb = INK
    if FONT:
        run.font.name = FONT


def _apply_takeaway_band_style(shape):
    """Match ``render_slide`` takeaway band fill and border on an existing shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = TAKEAWAY_FILL
    shape.line.fill.background()
    shape.shadow.inherit = False


def _fill_takeaway_text_frame(tf, takeaway):
    """Write takeaway label + body with theme colors (shared by render and stamp)."""
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = TAKEAWAY_TEXT_INSET
    tf.margin_right = TAKEAWAY_TEXT_INSET
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    _clear_paragraph_runs(para)
    lead = para.add_run()
    lead.text = "Takeaway:  "
    lead.font.bold = True
    lead.font.size = Pt(SIZES["takeaway"])
    lead.font.color.rgb = ACCENT
    if FONT:
        lead.font.name = FONT
    add_rich_text(para, takeaway, size=Pt(SIZES["takeaway"]), color=TAKEAWAY_TEXT)


def _chrome_fields(chrome):
    """Normalize brief or manifest chrome dict (``kicker`` vs ``role``)."""
    if not chrome:
        return {"role": None, "title": None, "takeaway": None}
    return {
        "role": chrome.get("role") or chrome.get("kicker"),
        "title": chrome.get("title"),
        "takeaway": chrome.get("takeaway"),
    }


def _add_slide_number_shape(slide, number, top=None):
    """Add top-right slide number textbox at standard chrome position."""
    chrome_top = CHROME_TOP if top is None else top
    num_w = Inches(1.2)
    nbox = slide.shapes.add_textbox(
        Emu(SLIDE_W - MARGIN - num_w), chrome_top, num_w, KICKER_H
    )
    _fill_slide_number_text_frame(nbox.text_frame, number)
    return nbox


def apply_slide_background(slide):
    """Set slide canvas fill from the active theme (SW render + HW stamp)."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND


def apply_slide_chrome(
    slide, chrome, *, number=None, top=None, mode="stamp", table_stripe=None
):
    """Apply themed slide chrome from one entry point.

    ``mode``:
      - ``stamp`` — find existing shapes (HW combine); restyle + set text. Returns missing fields.
      - ``create_header`` — add number/role/title shapes at ``top``; returns Emu below title.
      - ``create_takeaway`` — add bottom takeaway band from ``chrome['takeaway']``.
    """
    fields = _chrome_fields(chrome)
    missing = []

    if mode == "stamp":
        apply_slide_background(slide)

        if number is not None and SHOW_SLIDE_NUMBER:
            shape = find_slide_number_shape(slide)
            if shape is None:
                _add_slide_number_shape(slide, number)
            else:
                _fill_slide_number_text_frame(shape.text_frame, number)

        if fields["role"] and SHOW_KICKER:
            shape = find_kicker_shape(slide)
            if not shape:
                missing.append("kicker/role")
            else:
                _fill_kicker_text_frame(shape.text_frame, fields["role"])

        if fields["title"]:
            shape = find_title_shape(slide)
            if not shape:
                missing.append("title")
            else:
                _fill_title_text_frame(shape.text_frame, fields["title"])

        if fields["takeaway"] and SHOW_TAKEAWAY_BAND:
            shape = find_takeaway_shape(slide)
            if not shape:
                missing.append("takeaway")
            else:
                _apply_takeaway_band_style(shape)
                _fill_takeaway_text_frame(shape.text_frame, fields["takeaway"])
                _bring_shape_to_front(shape)

        recolor_transparent_text_boxes(slide)
        recolor_automatic_opaque_text_boxes(slide)
        recolor_automatic_diagram_lines(slide)
        recolor_table_shapes(slide, stripe=table_stripe)
        return missing

    if mode == "create_header":
        if top is None:
            fail("apply_slide_chrome(create_header) requires top=CHROME_TOP")

        header_row = False
        if number is not None and SHOW_SLIDE_NUMBER:
            _add_slide_number_shape(slide, number, top=top)
            header_row = True

        if fields["role"] and SHOW_KICKER:
            box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, KICKER_H)
            _fill_kicker_text_frame(box.text_frame, fields["role"])
            header_row = True

        if header_row:
            top = Emu(top + KICKER_H + GAP_AFTER_KICKER)

        if fields["title"]:
            title_h = _title_box_height(fields["title"])
            box = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, title_h)
            _fill_title_text_frame(box.text_frame, fields["title"])
            top = Emu(top + title_h + GAP_AFTER_TITLE)

        return top

    if mode == "create_takeaway":
        takeaway = fields["takeaway"]
        if not takeaway or not SHOW_TAKEAWAY_BAND:
            return None
        band = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN,
            Emu(SLIDE_H - TAKEAWAY_BOTTOM - TAKEAWAY_BAND_H),
            CONTENT_W,
            TAKEAWAY_BAND_H,
        )
        _apply_takeaway_band_style(band)
        _fill_takeaway_text_frame(band.text_frame, takeaway)
        return None

    fail(f'apply_slide_chrome: unknown mode "{mode}"')


def update_slide_number(slide, number):
    """Set the slide-number chrome text to ``number`` (combined-deck position)."""
    if not SHOW_SLIDE_NUMBER:
        return False
    shape = find_slide_number_shape(slide)
    if shape is None:
        return False
    _fill_slide_number_text_frame(shape.text_frame, number)
    return True


def find_kicker_shape(slide):
    """Top-left chrome label (Narrative Role / kicker band)."""
    top_max = int(SLIDE_H * 0.16)
    left_max = int(MARGIN + Inches(0.3))
    best = None
    best_top = 10**9
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top > top_max or shape.left > left_max:
            continue
        text = shape.text_frame.text.strip()
        if re.fullmatch(r"\d{1,3}", text):
            continue
        if shape.height > int(Inches(0.45)):
            continue
        if shape.top < best_top:
            best_top = shape.top
            best = shape
    return best


def _shape_key(shape):
    """Stable identity for a shape (``id()`` is not reliable across pptx iterators)."""
    return (shape.top, shape.left, shape.width, shape.height)


_SKIP_TRANSPARENT_RECOLOR_TYPES = {
    MSO_SHAPE_TYPE.PICTURE,
    MSO_SHAPE_TYPE.TABLE,
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.LINKED_PICTURE,
    MSO_SHAPE_TYPE.MEDIA,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
}


def _shape_has_solid_fill(shape):
    """True when the shape has an explicit solid fill (opaque label/badge)."""
    if not hasattr(shape, "fill"):
        return False
    try:
        fill_type = shape.fill.type
    except (AttributeError, TypeError):
        fill_type = None
    if fill_type == MSO_FILL_TYPE.SOLID:
        return True
    if fill_type is None:
        sp_pr = shape._element.spPr
        if sp_pr is not None:
            if sp_pr.find(f"{_DRAWING_TAG}solidFill") is not None:
                return True
    return False


def _chrome_shape_keys(slide):
    """Identity keys for shapes already handled as slide chrome."""
    keys = set()
    for finder in (
        find_slide_number_shape,
        find_kicker_shape,
        find_title_shape,
        find_takeaway_shape,
    ):
        found = finder(slide)
        if found:
            keys.add(_shape_key(found))
    return keys


# Theme slots PowerPoint treats as "Automatic" text (not a deliberate accent pick).
_AUTOMATIC_SCHEME_COLORS = frozenset({"tx1", "dk1", "lt1"})

# Light-mode authoring defaults for diagram strokes (PowerPoint black / near-black).
_AUTOMATIC_LINE_RGB = frozenset({"000000", "1A1A1A", "2A2420"})


def _shape_sp_pr(shape):
    if not hasattr(shape, "_element"):
        return None
    return getattr(shape._element, "spPr", None)


def _line_element(shape):
    sp_pr = _shape_sp_pr(shape)
    if sp_pr is None:
        return None
    return sp_pr.find(f"{_DRAWING_TAG}ln")


def _run_has_explicit_font_color(run):
    """True when the run has a deliberate font color (skip on opaque badge recolor)."""
    rpr = run._r.rPr
    if rpr is None:
        return False
    solid = rpr.find(f"{_DRAWING_TAG}solidFill")
    if solid is None:
        return False
    if solid.find(f"{_DRAWING_TAG}srgbClr") is not None:
        return True
    scheme = solid.find(f"{_DRAWING_TAG}schemeClr")
    if scheme is not None:
        return scheme.get("val") not in _AUTOMATIC_SCHEME_COLORS
    return True


def _set_text_frame_ink(tf):
    """Set all runs in a text frame to theme ink (diagram/callout labels)."""
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = INK
            if FONT:
                run.font.name = FONT


def _set_text_frame_automatic_ink(tf):
    """Set theme ink only on runs still using automatic / theme-default text color."""
    for para in tf.paragraphs:
        for run in para.runs:
            if _run_has_explicit_font_color(run):
                continue
            run.font.color.rgb = INK
            if FONT:
                run.font.name = FONT


def _recolor_transparent_shapes(shapes, skip):
    """Walk ``shapes`` (slide or group children); recurse into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _recolor_transparent_shapes(shape.shapes, skip)
            continue
        if shape.shape_type in _SKIP_TRANSPARENT_RECOLOR_TYPES:
            continue
        if not shape.has_text_frame:
            continue
        if _shape_key(shape) in skip:
            continue
        if _shape_has_solid_fill(shape):
            continue
        _set_text_frame_ink(shape.text_frame)


def recolor_transparent_text_boxes(slide):
    """After chrome stamp, set theme ink on transparent text boxes and callouts."""
    _recolor_transparent_shapes(slide.shapes, _chrome_shape_keys(slide))


def _recolor_automatic_opaque_shapes(shapes, skip):
    """Solid-fill text boxes/callouts: set ink on automatic-color runs only."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _recolor_automatic_opaque_shapes(shape.shapes, skip)
            continue
        if shape.shape_type in _SKIP_TRANSPARENT_RECOLOR_TYPES:
            continue
        if not shape.has_text_frame:
            continue
        if _shape_key(shape) in skip:
            continue
        if not _shape_has_solid_fill(shape):
            continue
        _set_text_frame_automatic_ink(shape.text_frame)


def recolor_automatic_opaque_text_boxes(slide):
    """After chrome stamp, set theme ink on automatic text in opaque callouts/badges."""
    _recolor_automatic_opaque_shapes(slide.shapes, _chrome_shape_keys(slide))


def _line_has_explicit_color(shape):
    """True when the shape stroke should not be recolored to theme ink."""
    if not hasattr(shape, "line"):
        return True
    ln = _line_element(shape)
    if ln is None:
        return True
    if ln.find(f"{_DRAWING_TAG}noFill") is not None:
        return True
    solid = ln.find(f"{_DRAWING_TAG}solidFill")
    if solid is None:
        return True  # no visible stroke — do not add one
    rgb = solid.find(f"{_DRAWING_TAG}srgbClr")
    if rgb is not None:
        val = (rgb.get("val") or "").upper()
        return val not in _AUTOMATIC_LINE_RGB
    scheme = solid.find(f"{_DRAWING_TAG}schemeClr")
    if scheme is not None:
        return scheme.get("val") not in _AUTOMATIC_SCHEME_COLORS
    return True


def _set_shape_line_ink(shape):
    """Set stroke color to theme ink; preserve any existing alpha on the stroke."""
    ln = _line_element(shape)
    alpha_val = None
    if ln is not None:
        solid = ln.find(f"{_DRAWING_TAG}solidFill")
        if solid is not None:
            rgb_el = solid.find(f"{_DRAWING_TAG}srgbClr")
            if rgb_el is not None:
                alpha = rgb_el.find(f"{_DRAWING_TAG}alpha")
                if alpha is not None:
                    alpha_val = alpha.get("val")
    shape.line.fill.solid()
    shape.line.color.rgb = INK
    if alpha_val is not None:
        ln = _line_element(shape)
        solid = ln.find(f"{_DRAWING_TAG}solidFill") if ln is not None else None
        rgb_el = solid.find(f"{_DRAWING_TAG}srgbClr") if solid is not None else None
        if rgb_el is not None and rgb_el.find(f"{_DRAWING_TAG}alpha") is None:
            rgb_el.append(
                parse_xml(f'<a:alpha xmlns:a="{_DRAWING_NS}" val="{alpha_val}"/>')
            )


def _recolor_automatic_line_shapes(shapes, skip):
    """Walk shapes; set ink on default-black / automatic diagram strokes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _recolor_automatic_line_shapes(shape.shapes, skip)
            continue
        if shape.shape_type in _SKIP_TRANSPARENT_RECOLOR_TYPES:
            continue
        if _shape_key(shape) in skip:
            continue
        if not hasattr(shape, "line"):
            continue
        if _line_has_explicit_color(shape):
            continue
        _set_shape_line_ink(shape)


def recolor_automatic_diagram_lines(slide):
    """After chrome stamp, set theme ink on automatic / default-black diagram strokes."""
    _recolor_automatic_line_shapes(slide.shapes, _chrome_shape_keys(slide))


def _normalize_table_stripe(stripe):
    """Return a validated stripe mode (row, column, none)."""
    if stripe is None:
        return TABLE_STRIPE
    mode = str(stripe).strip().lower()
    if mode not in VALID_TABLE_STRIPES:
        fail(f'Invalid table stripe mode "{stripe}" (expected row, column, or none).')
    return mode


def _set_cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def _set_cell_text_color(cell, rgb):
    tf = cell.text_frame
    for para in tf.paragraphs:
        if para.runs:
            for run in para.runs:
                run.font.color.rgb = rgb
        else:
            run = para.add_run()
            run.font.color.rgb = rgb


def _table_data_fill(band_index):
    """Alternate canvas background with one theme band color."""
    return BACKGROUND if band_index % 2 == 0 else TABLE_BAND_FILL


def _table_border_side_xml(side, hex_val, width_emu, *, ln_tag=None):
    """OOXML for one table/cell border side."""
    tag = ln_tag or side
    return (
        f"<a:{tag} w=\"{width_emu}\" cap=\"flat\" cmpd=\"sng\" algn=\"ctr\">"
        f"<a:solidFill><a:srgbClr val=\"{hex_val}\"/></a:solidFill>"
        f'<a:prstDash val="solid"/>'
        f"</a:{tag}>"
    )


def _set_cell_borders(cell, rgb, width_pt=TABLE_BORDER_WIDTH_PT):
    """Set all four cell borders to ``rgb`` (matches header frame color)."""
    tc_pr = cell._tc.get_or_add_tcPr()
    hex_val = str(rgb)
    width_emu = int(width_pt * 12700)
    for side in _TABLE_BORDER_TAG:
        tag = _TABLE_BORDER_TAG[side]
        existing = tc_pr.find(f"{_DRAWING_TAG}{tag}")
        if existing is not None:
            tc_pr.remove(existing)
        tc_pr.append(
            parse_xml(
                f'<a:{tag} xmlns:a="{_DRAWING_NS}" w="{width_emu}" cap="flat" '
                f'cmpd="sng" algn="ctr">'
                f'<a:solidFill><a:srgbClr val="{hex_val}"/></a:solidFill>'
                f'<a:prstDash val="solid"/></a:{tag}>'
            )
        )


def _set_table_level_borders(table, rgb, width_pt=TABLE_BORDER_WIDTH_PT):
    """Stamp explicit tblBorders on tblPr (overrides default theme tx1 grid)."""
    tbl_pr = table._tbl.tblPr
    existing = tbl_pr.find(f"{_DRAWING_TAG}tblBorders")
    if existing is not None:
        tbl_pr.remove(existing)
    hex_val = str(rgb)
    width_emu = int(width_pt * 12700)
    sides = "".join(
        _table_border_side_xml(side, hex_val, width_emu) for side in _TABLE_TBL_BORDER_SIDES
    )
    tbl_pr.append(
        parse_xml(f'<a:tblBorders xmlns:a="{_DRAWING_NS}">{sides}</a:tblBorders>')
    )


def _detach_table_style(table):
    """Remove embedded table style so tblBorders + cell tcPr control the grid."""
    tbl_pr = table._tbl.tblPr
    style_id = tbl_pr.find(f"{_DRAWING_TAG}tableStyleId")
    if style_id is not None:
        tbl_pr.remove(style_id)
    look = tbl_pr.find(f"{_DRAWING_TAG}tblLook")
    if look is None:
        tbl_pr.append(
            parse_xml(
                f'<a:tblLook xmlns:a="{_DRAWING_NS}" firstRow="0" lastRow="0" '
                f'firstCol="0" lastCol="0" bandRow="0" bandCol="0"/>'
            )
        )
    else:
        for attr in (
            "firstRow",
            "lastRow",
            "firstCol",
            "lastCol",
            "bandRow",
            "bandCol",
        ):
            look.set(attr, "0")


def recolor_table_shapes(slide, stripe=None):
    """Theme HW tables: row 0 + column 0 = header frame; data cells banded."""
    mode = _normalize_table_stripe(stripe)
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        _detach_table_style(table)
        _set_table_level_borders(table, TABLE_HEADER_FILL)
        rows = len(table.rows)
        cols = len(table.columns)
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                if r == 0 or c == 0:
                    _set_cell_fill(cell, TABLE_HEADER_FILL)
                    _set_cell_text_color(cell, TABLE_HEADER_TEXT)
                elif mode == "none":
                    _set_cell_fill(cell, BACKGROUND)
                    _set_cell_text_color(cell, TABLE_CELL_TEXT)
                elif mode == "column":
                    fill = _table_data_fill(c - 1)
                    _set_cell_fill(cell, fill)
                    _set_cell_text_color(cell, TABLE_CELL_TEXT)
                else:
                    fill = _table_data_fill(r - 1)
                    _set_cell_fill(cell, fill)
                    _set_cell_text_color(cell, TABLE_CELL_TEXT)
                _set_cell_borders(cell, TABLE_HEADER_FILL)


def find_title_shape(slide):
    """Title textbox below the kicker, above the body zone."""
    top_max = int(SLIDE_H * 0.28)
    left_max = int(MARGIN + Inches(0.3))
    skip_keys = set()
    for finder in (find_slide_number_shape, find_kicker_shape):
        found = finder(slide)
        if found:
            skip_keys.add(_shape_key(found))
    best = None
    best_top = -1
    for shape in slide.shapes:
        if _shape_key(shape) in skip_keys or not shape.has_text_frame:
            continue
        if shape.top > top_max or shape.left > left_max:
            continue
        if shape.text_frame.text.strip().lower().startswith("takeaway"):
            continue
        if shape.height < int(Inches(0.35)):
            continue
        if shape.top > best_top:
            best_top = shape.top
            best = shape
    return best


def _takeaway_band_top():
    return int(SLIDE_H - TAKEAWAY_BOTTOM - TAKEAWAY_BAND_H)


def _bring_shape_to_front(shape):
    """Move a shape to the top of the slide z-order (covers overlapping diagram art)."""
    el = shape._element
    parent = el.getparent()
    parent.remove(el)
    parent.append(el)


def find_takeaway_shape(slide):
    """Bottom takeaway band (rounded rectangle or wide textbox).

    HW source slides often contain diagram connector lines (zero-width AUTO_SHAPE)
    below the real band; prefer labeled bands and ignore narrow slivers.
    """
    band_top = _takeaway_band_top()
    bottom_min = int(SLIDE_H * 0.62)
    min_width = int(CONTENT_W * 0.5)
    min_height = int(TAKEAWAY_BAND_H * 0.4)
    band_tol = int(Inches(0.25))
    labeled = []
    bands = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top < bottom_min:
            continue
        if shape.width < min_width or shape.height < min_height:
            continue
        text = shape.text_frame.text.strip().lower()
        if "takeaway" in text:
            labeled.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            bands.append(shape)

    def _band_score(shape):
        in_band = 1 if abs(shape.top - band_top) <= band_tol else 0
        return (in_band, shape.width, -abs(shape.top - band_top))

    if labeled:
        return max(labeled, key=_band_score)
    if bands:
        return max(bands, key=lambda s: (s.width, -abs(s.top - band_top)))
    return None


def update_slide_kicker(slide, role):
    if not SHOW_KICKER:
        return False
    shape = find_kicker_shape(slide)
    if not shape:
        return False
    _fill_kicker_text_frame(shape.text_frame, role)
    return True


def update_slide_title(slide, title):
    shape = find_title_shape(slide)
    if not shape:
        return False
    _fill_title_text_frame(shape.text_frame, title)
    return True


def update_slide_takeaway(slide, takeaway):
    if not SHOW_TAKEAWAY_BAND:
        return False
    shape = find_takeaway_shape(slide)
    if not shape:
        return False
    _apply_takeaway_band_style(shape)
    _fill_takeaway_text_frame(shape.text_frame, takeaway)
    _bring_shape_to_front(shape)
    return True


def normalize_chrome_text(text):
    """Normalize manifest/Marp chrome strings for comparison."""
    if not text:
        return ""
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    t = t.replace("—", "-").replace("–", "-")
    t = re.sub(r"\s+", " ", t).strip()
    return t


def validate_markdown_chrome(tag, brief_slide, spine_row):
    """Return mismatch messages when Marp/brief chrome differs from the manifest."""
    mismatches = []
    exp_role = normalize_chrome_text(spine_row.get("role", "")).upper()
    got_role = normalize_chrome_text(brief_slide.get("kicker") or "").upper()
    if exp_role != got_role:
        mismatches.append(f"role/kicker: manifest {exp_role!r} vs Marp {got_role!r}")

    exp_title = normalize_chrome_text(spine_row.get("title", ""))
    got_title = normalize_chrome_text(brief_slide.get("title") or "")
    if exp_title != got_title:
        mismatches.append(f"title: manifest {exp_title!r} vs Marp {got_title!r}")

    exp_tw = normalize_chrome_text(spine_row.get("takeaway", ""))
    got_tw = normalize_chrome_text(brief_slide.get("takeaway") or "")
    if exp_tw != got_tw:
        mismatches.append(f"takeaway: manifest {exp_tw!r} vs Marp {got_tw!r}")

    return mismatches


def stamp_pptx_slide_chrome(slide, spine_row, number=None):
    """Write deck-manifest chrome onto PowerPoint slides. Returns missing fields."""
    return apply_slide_chrome(slide, spine_row, number=number, mode="stamp")


def render_slide(prs, blank, s, number, png_dir):
    """Render one parsed slide dict onto a new slide appended to ``prs``.

    ``number`` is the value shown top-right — the standalone index when called
    from build(), or the final combined position when called by the combine
    step. Returns the created slide.
    """
    slide = prs.slides.add_slide(blank)
    clear_slide_placeholders(slide)

    apply_slide_background(slide)

    top = apply_slide_chrome(
        slide, s, number=number, top=CHROME_TOP, mode="create_header"
    )

    # Vertical budget: content lives between `top` and the takeaway band.
    show_takeaway = bool(s["takeaway"]) and SHOW_TAKEAWAY_BAND
    takeaway_h = TAKEAWAY_BAND_H if show_takeaway else Inches(0.0)
    zone_top = top
    zone_bottom = Emu(SLIDE_H - TAKEAWAY_BOTTOM - takeaway_h)
    zone_h = zone_bottom - zone_top

    # Flow: diagram then body — matches Marp slide order (title → diagram → bullets).
    cursor = zone_top
    if s["diagram"]:
        png = Path(png_dir) / f"{s['diagram']}.png"
        if not png.exists():
            fail(f"Diagram PNG not found: {png}")
        pw, ph = png_size(png)
        ratio = pw / ph
        max_h = _diagram_max_height(zone_h, s.get("body"))
        if s.get("body"):
            body_cap = int(zone_h) - _body_reserve(s["body"]) - int(GAP_AFTER_DIAGRAM)
            max_h = min(max_h, max(0, body_cap))
        w = CONTENT_W
        h = int(CONTENT_W / ratio)
        if h > max_h:
            h = max_h
            w = int(max_h * ratio)
        left = Emu(MARGIN + (CONTENT_W - w) // 2)
        slide.shapes.add_picture(str(png), left, cursor, width=Emu(w), height=Emu(h))
        cursor = Emu(cursor + h + GAP_AFTER_DIAGRAM)

    if s["body"]:
        box = slide.shapes.add_textbox(MARGIN, cursor, CONTENT_W, Emu(zone_bottom - cursor))
        set_body(box.text_frame, s["body"])

    if show_takeaway:
        apply_slide_chrome(slide, s, mode="create_takeaway")

    if s["notes"]:
        slide.notes_slide.notes_text_frame.text = s["notes"]

    return slide


def build(brief_path, png_dir, out_path, theme_path=None, common_path=None):
    load_layout(common_path)
    if theme_path:
        load_theme(theme_path)
    text = Path(brief_path).read_text(encoding="utf-8")
    slides = parse_brief(text)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = pick_blank_layout(prs)

    for i, s in enumerate(slides, start=1):
        render_slide(prs, blank, s, i, png_dir)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path} ({len(slides)} slides).")


if __name__ == "__main__":
    if len(sys.argv) not in (4, 5):
        fail('Usage: brief-to-pptx.py "<brief.md>" "<png-dir>" "<out.pptx>" ["<palette.json>"]')
    theme = sys.argv[4] if len(sys.argv) == 5 else None
    build(sys.argv[1], sys.argv[2], sys.argv[3], theme)
