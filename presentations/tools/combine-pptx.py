#!/usr/bin/env python3
"""Assemble a combined .pptx by merging Markdown- and PowerPoint-sourced slides.

Reads a **deck note** (``<deck> Deck.md`` in the vault). The note's **deck manifest**
table defines final slide order and source tags (``S1…Sn`` from Markdown sources,
``H1…Hm`` from PowerPoint sources, etc.). Frontmatter lists those sources under
``deck_sources`` (each entry: ``file``, ``prefix``, ``type``, ``path``).

The first listed PowerPoint source *is* the output's starting point. Markdown slides
are generated into that file; PowerPoint-sourced slides get the themed slide
background and chrome stamped from the deck manifest; Markdown slide chrome is
validated against it. Then slide order is rearranged to match the manifest.

Usage:
  combine-pptx.py "<brief.md>" "<png-dir>" "<out.pptx>" "<deck-note.md>" ["<theme.json>"]
"""

import importlib.util
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation

# Import the sibling generator (hyphenated filename → load by path).
_B2P_PATH = Path(__file__).resolve().parent / "brief-to-pptx.py"
_spec = importlib.util.spec_from_file_location("brief_to_pptx", _B2P_PATH)
b2p = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(b2p)

TOKEN_RE = re.compile(r"^(S\d+|[A-Z]+\d+)$")
REPO_ROOT = Path(__file__).resolve().parent.parent.parent
BUILD_DIR = REPO_ROOT / "presentations" / "build"


def parse_frontmatter(path):
    """Return the YAML frontmatter dict from a markdown note."""
    text = Path(path).read_text(encoding="utf-8")
    if not text.startswith("---"):
        return {}
    end = text.find("\n---", 3)
    if end == -1:
        return {}
    try:
        data = yaml.safe_load(text[3:end]) or {}
    except yaml.YAMLError as exc:
        b2p.fail(f"Invalid YAML frontmatter in {path}: {exc}")
    return data if isinstance(data, dict) else {}


def normalize_table_stripe(value, *, context="table_stripe"):
    """Return ``row``, ``column``, or ``none``; fail on invalid values."""
    if value is None:
        return None
    mode = str(value).strip().lower()
    if mode in ("row", "column", "none"):
        return mode
    b2p.fail(f'Invalid {context} value {value!r} (expected row, column, or none).')


def resolve_table_stripe(manifest_path, tag):
    """Stripe mode for one HW slide tag from deck note ``table_stripe`` frontmatter."""
    raw = parse_frontmatter(manifest_path).get("table_stripe")
    if raw is None:
        return None
    if isinstance(raw, str):
        return normalize_table_stripe(raw)
    if isinstance(raw, dict):
        if tag in raw:
            return normalize_table_stripe(raw[tag], context=f"table_stripe.{tag}")
        if "default" in raw:
            return normalize_table_stripe(
                raw["default"], context="table_stripe.default"
            )
    b2p.fail(
        f"table_stripe in {manifest_path} must be a string or mapping "
        f"(default + per-tag keys like H1)."
    )


def _normalize_source(entry):
    """Return ``{file, prefix, type, path?}`` for one ``deck_sources`` entry."""
    if not isinstance(entry, dict):
        b2p.fail(
            "deck_sources entries must be mappings with file, prefix, type, and path, e.g.\n"
            "  - file: ORAM Software Slides\n"
            "    prefix: S\n"
            "    type: marp\n"
            "    path: vault/Notes/ORAM Software Slides.marp.md"
        )

    file_stem = (entry.get("file") or entry.get("stem") or "").strip()
    if not file_stem:
        b2p.fail("Each deck_sources entry needs file: …")

    prefix = (entry.get("prefix") or "").strip().upper()
    if not prefix or not re.fullmatch(r"[A-Z]+", prefix):
        b2p.fail(f"deck_sources prefix must be uppercase letters (got {prefix!r}).")

    path_val = (entry.get("path") or "").strip()
    type_val = (entry.get("type") or "").strip().lower()
    if not type_val and path_val:
        lower = path_val.lower()
        if lower.endswith(".marp.md"):
            type_val = "marp"
        elif lower.endswith(".pptx"):
            type_val = "pptx"
    if type_val not in ("marp", "pptx"):
        b2p.fail(
            f"deck_sources entry '{file_stem}' needs type: marp or pptx "
            f"(or a .marp.md / .pptx path)."
        )

    out = {"file": file_stem, "prefix": prefix, "type": type_val}
    if path_val:
        out["path"] = path_val
    return out


def parse_deck_sources(manifest_path):
    """Deck sources declared in manifest frontmatter."""
    fm = parse_frontmatter(manifest_path)
    raw = fm.get("deck_sources")
    if not raw:
        b2p.fail(
            f"No 'deck_sources' in frontmatter of {manifest_path} — declare inputs, e.g.\n"
            f"  deck_sources:\n"
            f"    - file: ORAM Software Slides\n"
            f"      prefix: S\n"
            f"      type: marp\n"
            f"      path: vault/Notes/ORAM Software Slides.marp.md\n"
            f"    - file: ORAM Hardware Slides\n"
            f"      prefix: H\n"
            f"      type: pptx\n"
            f"      path: presentations/build/ORAM Hardware Slides.pptx"
        )
    if isinstance(raw, str):
        raw = [raw]
    if not isinstance(raw, list) or not raw:
        b2p.fail(f"deck_sources in {manifest_path} must be a non-empty list.")

    sources = [_normalize_source(entry) for entry in raw]
    prefixes = [s["prefix"] for s in sources]
    if len(set(prefixes)) != len(prefixes):
        b2p.fail(f"Duplicate deck_sources prefixes in {manifest_path}: {prefixes}")
    return sources


def resolve_source_path(path_str, default_repo_relative=None):
    """Resolve a deck source path.

    - Absolute paths and ``~`` are used as-is (after expanduser).
    - Relative paths are resolved from the repository root.
    - When ``path_str`` is empty, ``default_repo_relative`` is used the same way.
    """
    raw = (path_str or "").strip() or (default_repo_relative or "").strip()
    if not raw:
        return None
    p = Path(raw).expanduser()
    if not p.is_absolute():
        p = REPO_ROOT / p
    return p


def resolve_sources(manifest_path):
    """Resolve every ``deck_sources`` entry to an absolute path."""
    resolved = []
    for src in parse_deck_sources(manifest_path):
        if src["type"] == "marp":
            default = f"vault/Notes/{src['file']}.marp.md"
        else:
            default = f"presentations/build/{src['file']}.pptx"
        path = resolve_source_path(src.get("path"), default)
        if not path.is_file():
            b2p.fail(
                f"Source not found: {path} "
                f"(deck_sources → {src['file']}, type={src['type']} in {manifest_path})"
            )
        resolved.append({**src, "path": path})
    return resolved


def parse_tag_prefix(tag):
    """Return the uppercase prefix from ``S3`` / ``H1``-style tags, else ``None``."""
    m = re.match(r"^([A-Z]+)(\d+)$", tag)
    return m.group(1) if m else None


def _manifest_header_indices(cells):
    """Column map if this header row starts a deck-manifest table, else ``None``.

    Skips summary tables (e.g. Half | Source | … listing file paths) that have
    ``Source`` but not ``Title`` or ``Narrative Role``.
    """
    cols = {}
    for i, c in enumerate(cells):
        key = c.lower().strip()
        if key == "source":
            cols["source"] = i
        elif "narrative role" in key or key == "role":
            cols["role"] = i
        elif key == "title":
            cols["title"] = i
        elif key == "takeaway":
            cols["takeaway"] = i
    if "source" not in cols:
        return None
    if "title" not in cols and "role" not in cols:
        return None
    return cols


def _normalize_source_tag(cell):
    return cell.strip().strip("`").strip()


def parse_spine_metadata(path):
    """Return ``{source_tag: {role, title, takeaway}}`` from the deck manifest table."""
    lines = Path(path).read_text(encoding="utf-8").splitlines()
    best = {}
    cols = None
    meta = {}
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|"):
            if cols is not None:
                if len(meta) > len(best):
                    best = meta
                cols = None
                meta = {}
            continue
        cells = [c.strip() for c in stripped.strip("|").split("|")]
        if all(re.fullmatch(r":?-+:?", c) for c in cells if c):
            continue
        if cols is None:
            cols = _manifest_header_indices(cells)
            continue
        if cols["source"] >= len(cells):
            continue
        tag = _normalize_source_tag(cells[cols["source"]])
        if not TOKEN_RE.match(tag):
            continue
        meta[tag] = {
            "role": cells[cols["role"]].strip() if "role" in cols and cols["role"] < len(cells) else "",
            "title": cells[cols["title"]].strip() if "title" in cols and cols["title"] < len(cells) else "",
            "takeaway": cells[cols["takeaway"]].strip() if "takeaway" in cols and cols["takeaway"] < len(cells) else "",
        }
    if cols is not None and len(meta) > len(best):
        best = meta
    if not best:
        b2p.fail(f"No deck manifest metadata found in {path} (table with Source column).")
    return best


def parse_manifest(path):
    """Return the ordered list of source tags from the deck manifest table."""
    lines = Path(path).read_text(encoding="utf-8").splitlines()
    best = []
    cols = None
    tags = []
    saw_manifest_header = False
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|"):
            if cols is not None:
                if len(tags) > len(best):
                    best = tags
                cols = None
                tags = []
            continue
        cells = [c.strip() for c in stripped.strip("|").split("|")]
        if all(re.fullmatch(r":?-+:?", c) for c in cells if c):
            continue
        if cols is None:
            cols = _manifest_header_indices(cells)
            if cols is not None:
                saw_manifest_header = True
            continue
        if cols["source"] >= len(cells):
            continue
        tag = _normalize_source_tag(cells[cols["source"]])
        if TOKEN_RE.match(tag):
            tags.append(tag)
    if cols is not None and len(tags) > len(best):
        best = tags
    if not saw_manifest_header:
        b2p.fail(f"No table with a 'Source' column found in {path}")
    if not best:
        b2p.fail(f"Found a 'Source' column but no source tags in {path}")
    return best


def combine(brief_path, png_dir, out_path, manifest_path, theme_path=None):
    b2p.load_layout()
    if theme_path:
        b2p.load_theme(theme_path)

    all_sources = resolve_sources(manifest_path)
    md_sources = [s for s in all_sources if s["type"] == "marp"]
    pptx_sources = [s for s in all_sources if s["type"] == "pptx"]

    if len(md_sources) > 1:
        b2p.fail(
            "Multiple marp sources are declared but combine currently supports "
            "one Marp deck — use a single entry until multi-source merge is implemented."
        )
    if len(pptx_sources) > 1:
        b2p.fail(
            "Multiple pptx sources are declared but combine currently supports "
            "one external deck — use a single entry until multi-source merge "
            "is implemented."
        )
    if not md_sources:
        b2p.fail(f"No marp source in deck_sources (manifest {manifest_path}).")
    if not pptx_sources:
        b2p.fail(f"No pptx source in deck_sources (manifest {manifest_path}).")

    md = md_sources[0]
    md_prefix = md["prefix"]
    base = pptx_sources[0]
    base_path = base["path"]
    base_stem = base["file"]
    ext_prefix = base["prefix"]

    md_prefixes = {s["prefix"] for s in md_sources}
    pptx_prefixes = {s["prefix"] for s in pptx_sources}
    all_prefixes = md_prefixes | pptx_prefixes

    slides = b2p.parse_brief(Path(brief_path).read_text(encoding="utf-8"))
    order = parse_manifest(manifest_path)
    spine = parse_spine_metadata(manifest_path)

    md_tags = [t for t in order if parse_tag_prefix(t) in md_prefixes]
    ext_tags = [t for t in order if parse_tag_prefix(t) in pptx_prefixes]

    for tag in order:
        prefix = parse_tag_prefix(tag)
        if prefix and prefix not in all_prefixes:
            b2p.fail(
                f"Deck manifest tag {tag} uses prefix {prefix} but no matching entry "
                f"in deck_sources (declared prefixes: {sorted(all_prefixes)})."
            )
        if tag not in spine:
            b2p.fail(f"Source tag {tag} missing from deck manifest in {manifest_path}.")

    chrome_errors = []
    for idx, slide in enumerate(slides, start=1):
        tag = f"{md_prefix}{idx}"
        mismatches = b2p.validate_markdown_chrome(tag, slide, spine[tag])
        if mismatches:
            chrome_errors.append(f"{tag}: " + "; ".join(mismatches))
    if chrome_errors:
        b2p.fail(
            "Markdown slide chrome does not match the deck manifest "
            f"(edit [[{md['file']}.marp]] or the deck manifest table):\n"
            + "\n".join(chrome_errors)
        )

    if len(md_tags) != len(slides):
        b2p.fail(
            f"Deck manifest lists {len(md_tags)} markdown slides ({md_prefix}#) but the "
            f"brief has {len(slides)}. Reconcile {manifest_path} with "
            f"{md['path'].name}."
        )

    prs = Presentation(base_path)
    prs.slide_width = b2p.SLIDE_W
    prs.slide_height = b2p.SLIDE_H
    blank = b2p.pick_blank_layout(prs)

    sld_id_lst = prs.slides._sldIdLst
    base_ids = list(sld_id_lst)
    n_ext = len(base_ids)
    if len(ext_tags) != n_ext:
        b2p.fail(
            f"Deck manifest lists {len(ext_tags)} PowerPoint slides ({ext_prefix}#) "
            f"but {base_path} has {n_ext} slides."
        )

    for idx, s in enumerate(slides, start=1):
        tag = f"{md_prefix}{idx}"
        final_pos = order.index(tag) + 1
        b2p.render_slide(prs, blank, s, final_pos, png_dir)

    missing_chrome = []
    for k in range(n_ext):
        tag = f"{ext_prefix}{k + 1}"
        final_pos = order.index(tag) + 1
        slide = prs.slides[k]
        row = spine[tag]
        table_stripe = resolve_table_stripe(manifest_path, tag)
        fields = b2p.apply_slide_chrome(
            slide, row, number=final_pos, mode="stamp", table_stripe=table_stripe
        )
        for field in fields:
            missing_chrome.append(f"{tag} ({field})")
    if missing_chrome:
        b2p.fail(
            f"Could not update slide chrome on PowerPoint slide(s): "
            f"{', '.join(missing_chrome)}. Each PowerPoint source slide needs chrome "
            f"textboxes matching the generated layout (number, kicker, title, takeaway)."
        )

    md_ids = list(sld_id_lst)[n_ext:]

    tag_to_el = {f"{ext_prefix}{k + 1}": el for k, el in enumerate(base_ids)}
    tag_to_el.update({f"{md_prefix}{k + 1}": el for k, el in enumerate(md_ids)})

    for el in list(sld_id_lst):
        sld_id_lst.remove(el)
    for tag in order:
        if tag not in tag_to_el:
            b2p.fail(f"Deck manifest source tag {tag} has no matching slide.")
        sld_id_lst.append(tag_to_el[tag])

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(
        f"Wrote {out_path} ({len(order)} slides: {n_ext} PowerPoint + {len(slides)} "
        f"markdown; base={base_stem}.pptx [{ext_prefix}#], marp={md['file']}.marp.md "
        f"[{md_prefix}#], ordered per {Path(manifest_path).name})."
    )


if __name__ == "__main__":
    if len(sys.argv) not in (5, 6):
        b2p.fail(
            'Usage: combine-pptx.py "<brief.md>" "<png-dir>" "<out.pptx>" '
            '"<manifest.md>" ["<theme.json>"]'
        )
    theme = sys.argv[5] if len(sys.argv) == 6 else None
    combine(sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4], theme)
